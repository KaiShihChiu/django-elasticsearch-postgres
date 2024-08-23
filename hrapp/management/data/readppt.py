from pptx import Presentation
from PIL import Image
import io
import os
import re


class PptExtractor:
    def __init__(self, file_path, output_dir):
        self.file_path = file_path
        self.output_dir = output_dir
        self.presentation = Presentation(file_path)
        self.sections = [
            "技術名稱", "預期用途", "技術架構", "競爭優勢", "應用範圍", "技術合作夥伴",
            "研究成果和獎項", "專利", "專利佈局情況", "技術聯絡人", "關鍵字", "技術成熟度"
        ]
        self.slides_content = self.extract_text_per_slide()

    def extract_text_per_slide(self):
        """Extracts text content from each slide in the presentation."""
        slides_content = []
        for i, slide in enumerate(self.presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            slides_content.append(f"Slide {i + 1}:\n" + "\n".join(slide_text))
        return slides_content

    @staticmethod
    def clean_text(text):
        """Cleans the text by removing multiple consecutive newlines and leading punctuation."""
        cleaned_text = re.sub(r'\n+', '\n', text).strip()  # Remove multiple newlines
        cleaned_text = re.sub(r'^[\s:：,，.。;；!！?？]+', '', cleaned_text)  # Remove leading punctuation
        return cleaned_text

    def extract_sections_and_images_list(self):
        """Extracts defined sections and images from each slide's content."""
        full_content = []

        for slide_num, slide_text in enumerate(self.slides_content, start=1):
            content = []

            # Identify the order of sections in the actual slide
            present_sections = []
            for sec in self.sections:
                match = re.search(f"\\s*{re.escape(sec)}\\s*", slide_text)
                if match:
                    present_sections.append((match.start(), sec))

            # Sort sections by their position in the text
            present_sections.sort()

            # Extract content between present sections
            for idx, (start_pos, sec) in enumerate(present_sections):
                next_sec_start = present_sections[idx + 1][0] if idx + 1 < len(present_sections) else None
                if next_sec_start:
                    section_content = slide_text[start_pos:next_sec_start].strip()
                else:
                    section_content = slide_text[start_pos:].strip()

                # Extract the text after the section title
                pattern = re.compile(f"\\s*{re.escape(sec)}\\s*(.*)", re.S)
                match = pattern.search(section_content)
                if match:
                    cleaned_content = self.clean_text(match.group(1))
                    content.append((sec, cleaned_content))
                else:
                    content.append((sec, "未找到匹配的內容。"))

            # Extract images from the slide
            images = self.extract_images_from_slide(slide_num)
            if images:
                content.append(("技術成品圖", images))

            full_content.append(content)

        return full_content
    
    def extract_sections_and_images(self):
        """Extracts defined sections and images from each slide's content."""
        full_content = []
        titles = []

        for slide_num, slide_text in enumerate(self.slides_content, start=1):
            content = []
            title = None

            # Identify the order of sections in the actual slide
            present_sections = []
            for sec in self.sections:
                match = re.search(f"\\s*{re.escape(sec)}\\s*", slide_text)
                if match:
                    present_sections.append((match.start(), sec))

            # Sort sections by their position in the text
            present_sections.sort()

            # Extract content between present sections
            for idx, (start_pos, sec) in enumerate(present_sections):
                next_sec_start = present_sections[idx + 1][0] if idx + 1 < len(present_sections) else None
                if next_sec_start:
                    section_content = slide_text[start_pos:next_sec_start].strip()
                else:
                    section_content = slide_text[start_pos:].strip()

                # Extract the text after the section title
                pattern = re.compile(f"\\s*{re.escape(sec)}\\s*(.*)", re.S)
                match = pattern.search(section_content)
                if match:
                    cleaned_content = self.clean_text(match.group(1))
                    # 如果当前的 section 是 "技術名稱"，将内容保存到 title 变量中
                    if sec == "技術名稱":
                        title = cleaned_content
                    else:
                        content.append(f"{sec}\n{cleaned_content}")
                else:
                    content.append(f"{sec}\n未找到匹配的內容。")

            # Join all sections content into a single description string, separated by \n\n
            description = "\n\n".join(content)
            self.extract_images_from_slide(slide_num)

            # Append the title and description to the lists
            titles.append(title)
            full_content.append(description)

        return titles, full_content

    def extract_images_from_slide(self, slide_num):
        """Extracts images from a specific slide and returns a list of image paths."""
        images = []
        slide = self.presentation.slides[slide_num - 1]  # Slide index is 0-based in presentation
        image_index = 1

        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = io.BytesIO(image.blob)
                img = Image.open(image_bytes)
                # the ppt file name should be included for image name
                # img_path = os.path.join(self.output_dir, f"slide_{slide_num}_image_{image_index}.png")
                img_path = os.path.join(self.output_dir, f"slide_{slide_num}_image_{image_index}.png")
                img.save(img_path)
                images.append(img_path)
                print(f"Saved image {image_index} from slide {slide_num} to {img_path}")
                image_index += 1

        return images
    
    # def save_to_patent_model(self, extracted_content):
    #     for content in extracted_content:
    #         title = None
    #         description = ""
    #         image_paths = []

    #         for item in content:
    #             section, value = item
    #             if section == "技術名稱":
    #                 title = value
    #             elif section == "技術成品圖":
    #                 image_paths.extend(value)
    #             else:
    #                 description += f"{section}: {value}\n"

    #         if title:
    #             # 保存数据到数据库
    #             patent = Patent.objects.create(title=title, description=description)
    #             for image_path in image_paths:
    #                 # 假设您有其他字段来保存图片路径
    #                 patent.image_path = image_path
    #                 patent.save()



# if __name__ == "__main__":
#     # extract the local directory path of the pptx file and the output directory path for the extracted images
#     local_dir = os.path.dirname(os.path.abspath(__file__))
#     print(local_dir)

#     # find the parent folder of the local directory
#     # parent_dir = os.path.dirname(local_dir)
#     # print(parent_dir)

#     # # Example usage:
#     file_path = os.path.join(local_dir, "reference.pptx")
#     output_dir = os.path.join(local_dir, "image")

#     extractor = PptExtractor(file_path, output_dir)
#     extracted_content = extractor.extract_sections_and_images()
#     # # extractor.save_to_patent_model(extracted_content)
    
#     for i, content in enumerate(extracted_content, 1):
#         print(f"Slide {i} content:")
#         print(content)  # 這裡的 content 已經是一個完整的 description 字段的內容
#         print("\n" + "="*40 + "\n")

    # Print extracted content (matched for list)
    # for i, content in enumerate(extracted_content, 1):
    #     print(f"Slide {i} content:")
    #     for item in content:
    #         section, value = item
    #         if isinstance(value, list):
    #             print(f"{section}:")
    #             for image_path in value:
    #                 print(f"  Image: {image_path}")
    #         else:
    #             print(f"{section}: {value}")
    #     print("\n" + "="*40 + "\n")
